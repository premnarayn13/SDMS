"""
PowerPoint (PPTX) power tools
Uses python-pptx for PPTX processing. Async-friendly executor interface.
"""
import io
import logging
from typing import Tuple, Optional

from pptx import Presentation

logger = logging.getLogger(__name__)


class PptPowerToolsExecutor:
    def __init__(self):
        pass

    async def extract_text(self, content: bytes) -> Tuple[Optional[str], Optional[str]]:
        try:
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts), None
        except Exception as e:
            logger.error("PPT extract_text error: %s", e)
            return None, str(e)

    async def split_slides(self, content: bytes) -> Tuple[Optional[list[bytes]], Optional[str]]:
        try:
            prs = Presentation(io.BytesIO(content))
            slides_bytes = []
            for i, slide in enumerate(prs.slides):
                new_prs = Presentation()
                # copy slide using XML hack - not perfect, placeholder
                new_prs.slides._sldIdLst.append(slide._element)
                buf = io.BytesIO()
                new_prs.save(buf)
                slides_bytes.append(buf.getvalue())
            return slides_bytes, None
        except Exception as e:
            logger.error("PPT split_slides error: %s", e)
            return None, str(e)

    async def merge_presentations(self, contents: list[bytes]) -> Tuple[Optional[bytes], Optional[str]]:
        try:
            # Naive merge: append slides from other presentations
            base = Presentation(io.BytesIO(contents[0]))
            for other in contents[1:]:
                prs = Presentation(io.BytesIO(other))
                for slide in prs.slides:
                    # This is non-trivial; placeholder says not implemented fully
                    base.slides.add_slide(slide.slide_layout)
            buf = io.BytesIO()
            base.save(buf)
            return buf.getvalue(), None
        except Exception as e:
            logger.error("PPT merge error: %s", e)
            return None, str(e)

    async def add_watermark(self, content: bytes, watermark_text: str) -> Tuple[Optional[bytes], Optional[str]]:
        # Placeholder: adding watermarks to slides requires drawing shapes
        return None, "Watermarking PPTX not implemented (requires slide graphics operations)."


_executor: PptPowerToolsExecutor = None


def get_ppt_executor() -> PptPowerToolsExecutor:
    global _executor
    if _executor is None:
        _executor = PptPowerToolsExecutor()
    return _executor
