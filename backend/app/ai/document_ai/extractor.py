"""
Document Text Extractor
Multi-format extractor leveraging Google Cloud Document AI for PDFs/Images and native Python libraries for editable document formats.
"""
import os
import io
import re
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path
from dotenv import load_dotenv

# Load backend/.env file
env_path = Path(__file__).resolve().parent.parent.parent.parent / ".env"
load_dotenv(dotenv_path=env_path)

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """Document Extractor class routing files to DocAI or local parsers."""

    def __init__(self):
        # Refresh GCP Document AI Credentials & Settings
        self.gcp_project_id = os.environ.get("GOOGLE_CLOUD_PROJECT") or os.environ.get("GCP_PROJECT_ID")
        self.gcp_location = os.environ.get("GOOGLE_DOCUMENT_AI_LOCATION", "us")
        self.gcp_processor_id = os.environ.get("GOOGLE_DOCUMENT_AI_PROCESSOR_ID")
        
        creds_env = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")
        if creds_env and not os.path.isabs(creds_env):
            backend_root = Path(__file__).resolve().parent.parent.parent.parent
            resolved_creds = backend_root / creds_env
            if resolved_creds.exists():
                os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = str(resolved_creds)
            elif (backend_root / "Credentials" / "document-ai.json.json").exists():
                os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = str(backend_root / "Credentials" / "document-ai.json.json")

    def extract_document(self, file_path: str, file_type: Optional[str] = None) -> Dict[str, Any]:
        """
        Main entry point for document extraction.
        Routes file based on extension/type.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Document file not found at path: {file_path}")

        ext = (file_type or path.suffix.lstrip(".")).lower()

        # Handle PDFs and Images via Google Document AI (with local fallback)
        if ext in ["pdf", "png", "jpg", "jpeg", "tiff", "tif"]:
            return self._extract_pdf_or_image(file_path, ext)
        elif ext in ["docx", "doc"]:
            return self._extract_docx(file_path)
        elif ext in ["xlsx", "xls", "csv"]:
            return self._extract_spreadsheet(file_path, ext)
        elif ext in ["pptx", "ppt"]:
            return self._extract_pptx(file_path)
        elif ext in ["txt", "md", "markdown", "json", "xml", "html"]:
            return self._extract_plain_text(file_path)
        else:
            # Fallback to plain text read attempt
            return self._extract_plain_text(file_path)

    def _extract_pdf_or_image(self, file_path: str, ext: str) -> Dict[str, Any]:
        """Attempts Google Document AI extraction first; falls back to pdfplumber/fitz/pypdf."""
        # Try Google Cloud Document AI if configured
        if self.gcp_project_id and self.gcp_processor_id:
            try:
                return self._extract_with_google_doc_ai(file_path, ext)
            except Exception as e:
                logger.warning(f"Google Document AI extraction failed/unavailable: {e}. Falling back to local PDF parser.")

        # Local Fallback
        return self._extract_pdf_local(file_path)

    def _extract_with_google_doc_ai(self, file_path: str, ext: str) -> Dict[str, Any]:
        """Uses google-cloud-documentai Python SDK."""
        from google.cloud import documentai_v1 as documentai

        client = documentai.DocumentProcessorServiceClient()
        name = client.processor_path(self.gcp_project_id, self.gcp_location, self.gcp_processor_id)

        mime_type = "application/pdf" if ext == "pdf" else f"image/{ext}"
        if ext in ["jpg", "jpeg"]:
            mime_type = "image/jpeg"

        with open(file_path, "rb") as image:
            image_content = image.read()

        raw_document = documentai.RawDocument(content=image_content, mime_type=mime_type)
        request = documentai.ProcessRequest(name=name, raw_document=raw_document)
        result = client.process_document(request=request)
        document = result.document

        pages = []
        for i, page in enumerate(document.pages, start=1):
            page_text = ""
            for paragraph in page.paragraphs:
                page_text += self._get_text_from_layout(paragraph.layout, document.text) + "\n"
            pages.append({"page_number": i, "text": page_text.strip()})

        full_text = document.text or "\n\n".join([p["text"] for p in pages])
        return {
            "raw_text": full_text,
            "pages": pages if pages else [{"page_number": 1, "text": full_text}],
            "tables": [],
            "forms": [],
            "extractor_used": "Google Document AI",
            "file_size": os.path.getsize(file_path)
        }

    def _get_text_from_layout(self, layout, full_text):
        """Utility for Document AI bounding box text reconstruction."""
        text = ""
        for segment in layout.text_anchor.text_segments:
            start_index = int(segment.start_index)
            end_index = int(segment.end_index)
            text += full_text[start_index:end_index]
        return text

    def _extract_pdf_local(self, file_path: str) -> Dict[str, Any]:
        """Local PDF extractor with PyMuPDF & OCR fallback for scanned image PDFs."""
        pages = []
        full_text_list = []

        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for idx, page in enumerate(pdf.pages, start=1):
                    txt = page.extract_text() or ""
                    pages.append({"page_number": idx, "text": txt.strip()})
                    full_text_list.append(txt.strip())
        except Exception:
            try:
                from pypdf import PdfReader
                reader = PdfReader(file_path)
                for idx, page in enumerate(reader.pages, start=1):
                    txt = page.extract_text() or ""
                    pages.append({"page_number": idx, "text": txt.strip()})
                    full_text_list.append(txt.strip())
            except Exception as e:
                logger.error(f"Error reading PDF text streams: {e}")

        # Check if text is sparse or empty (Scanned PDF Document)
        combined_text = "".join(full_text_list).strip()
        if len(combined_text) < 30:
            logger.info("PDF text stream empty or sparse. Running Scanned PDF Image OCR...")
            try:
                import fitz  # PyMuPDF
                from PIL import Image
                import io
                import base64

                doc = fitz.open(file_path)
                pages = []
                full_text_list = []

                for idx, page in enumerate(doc, start=1):
                    pix = page.get_pixmap(dpi=150)
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    
                    ocr_text = ""
                    # 1. Try PyTesseract OCR
                    try:
                        import pytesseract
                        ocr_text = pytesseract.image_to_string(img).strip()
                    except Exception as ocr_err:
                        logger.debug(f"PyTesseract OCR bypass: {ocr_err}")

                    # 2. Fallback to Groq Vision OCR if Tesseract produces no text
                    if not ocr_text and self.groq_api_key:
                        try:
                            from groq import Groq
                            buffered = io.BytesIO()
                            img.save(buffered, format="PNG")
                            img_b64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                            
                            client = Groq(api_key=self.groq_api_key)
                            vision_resp = client.chat.completions.create(
                                model="llama-3.2-11b-vision-preview",
                                messages=[
                                    {
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": "Transcribe all text from this scanned image page accurately. Output only the verbatim text found."},
                                            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}
                                        ]
                                    }
                                ],
                                max_tokens=1000
                            )
                            ocr_text = vision_resp.choices[0].message.content.strip()
                        except Exception as vis_err:
                            logger.warning(f"Groq Vision OCR error on page {idx}: {vis_err}")

                    if not ocr_text:
                        ocr_text = f"[Scanned Image Page {idx} - No text recognized]"

                    pages.append({"page_number": idx, "text": ocr_text})
                    full_text_list.append(ocr_text)

                doc.close()
            except Exception as scanned_err:
                logger.error(f"Error in scanned PDF OCR pipeline: {scanned_err}")

        full_text = "\n\n--- Page Break ---\n\n".join(full_text_list)
        return {
            "raw_text": full_text,
            "pages": pages if pages else [{"page_number": 1, "text": full_text}],
            "tables": [],
            "forms": [],
            "extractor_used": "Scanned PDF OCR Engine (PyMuPDF + Tesseract/Vision)",
            "file_size": os.path.getsize(file_path)
        }

    def _extract_docx(self, file_path: str) -> Dict[str, Any]:
        """Extracts text from DOCX using python-docx."""
        try:
            import docx
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n".join(paragraphs)

            # Basic page estimation (every 500 words is ~1 page)
            words = full_text.split()
            page_count = max(1, len(words) // 450)
            pages = []
            words_per_page = max(1, len(words) // page_count) if words else 1
            
            for i in range(page_count):
                chunk = " ".join(words[i*words_per_page : (i+1)*words_per_page])
                pages.append({"page_number": i+1, "text": chunk})

            return {
                "raw_text": full_text,
                "pages": pages,
                "tables": [],
                "forms": [],
                "extractor_used": "python-docx",
                "file_size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"Error reading DOCX: {e}")
            return self._extract_plain_text(file_path)

    def _extract_spreadsheet(self, file_path: str, ext: str) -> Dict[str, Any]:
        """Extracts spreadsheet content using pandas / openpyxl."""
        try:
            import pandas as pd
            if ext == "csv":
                df = pd.read_csv(file_path)
                full_text = df.to_string()
            else:
                excel_file = pd.ExcelFile(file_path)
                sheets_text = []
                for sheet in excel_file.sheet_names:
                    df = pd.read_excel(excel_file, sheet_name=sheet)
                    sheets_text.append(f"=== Sheet: {sheet} ===\n" + df.to_string())
                full_text = "\n\n".join(sheets_text)

            return {
                "raw_text": full_text,
                "pages": [{"page_number": 1, "text": full_text}],
                "tables": [],
                "forms": [],
                "extractor_used": "pandas/openpyxl",
                "file_size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"Error reading spreadsheet: {e}")
            return self._extract_plain_text(file_path)

    def _extract_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extracts presentation content using python-pptx."""
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            pages = []
            slide_texts = []

            for idx, slide in enumerate(prs.slides, start=1):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
                slide_content = "\n".join(text_runs)
                pages.append({"page_number": idx, "text": slide_content})
                slide_texts.append(f"--- Slide {idx} ---\n" + slide_content)

            full_text = "\n\n".join(slide_texts)
            return {
                "raw_text": full_text,
                "pages": pages,
                "tables": [],
                "forms": [],
                "extractor_used": "python-pptx",
                "file_size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"Error reading PPTX: {e}")
            return self._extract_plain_text(file_path)

    def _extract_plain_text(self, file_path: str) -> Dict[str, Any]:
        """Fallback plain text reader."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                full_text = f.read()

            words = full_text.split()
            page_count = max(1, len(words) // 450)
            pages = []
            words_per_page = max(1, len(words) // page_count) if words else 1
            for i in range(page_count):
                chunk = " ".join(words[i*words_per_page : (i+1)*words_per_page])
                pages.append({"page_number": i+1, "text": chunk})

            return {
                "raw_text": full_text,
                "pages": pages if pages else [{"page_number": 1, "text": full_text}],
                "tables": [],
                "forms": [],
                "extractor_used": "native-text-reader",
                "file_size": os.path.getsize(file_path)
            }
        except Exception as e:
            logger.error(f"Error reading plain text file: {e}")
            return {
                "raw_text": f"Error reading document file: {str(e)}",
                "pages": [{"page_number": 1, "text": ""}],
                "tables": [],
                "forms": [],
                "extractor_used": "error-fallback",
                "file_size": 0
            }


extractor_instance = DocumentExtractor()
